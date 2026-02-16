#!/usr/bin/env python3
"""
Convert Executive AI Briefing HTML slides to PowerPoint format.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from bs4 import BeautifulSoup
import re

# Color palette (from CSS variables)
COLORS = {
    'primary': RGBColor(0x6B, 0x8E, 0x23),      # Olive green
    'coral': RGBColor(0xFF, 0x6B, 0x6B),
    'blue': RGBColor(0x45, 0xAE, 0xEE),
    'purple': RGBColor(0xA8, 0x6D, 0xD8),
    'teal': RGBColor(0x20, 0xC9, 0x97),
    'gold': RGBColor(0xF4, 0xC5, 0x42),
    'dark': RGBColor(0x1A, 0x1A, 0x2E),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'text': RGBColor(0x2D, 0x3A, 0x4A),
}


def create_presentation():
    """Create a new presentation with 16:9 aspect ratio."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_title_slide(prs, title, subtitle=None, logo=None):
    """Add a title slide."""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['dark']
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = COLORS['primary']
        p.alignment = PP_ALIGN.CENTER

    # Logo
    if logo:
        logo_box = slide.shapes.add_textbox(Inches(5.5), Inches(5.5), Inches(2.333), Inches(0.6))
        tf = logo_box.text_frame
        p = tf.paragraphs[0]
        p.text = logo
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_section_title_slide(prs, title, subtitle=None):
    """Add a section title slide with accent color background."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['primary']
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, content_items=None, subtitle=None):
    """Add a content slide with title and bullet points."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0xFA, 0xF9, 0xF6)
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['text']

    # Subtitle
    y_offset = 1.5
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.4), Inches(11.833), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['primary']
        y_offset = 2.0

    # Content
    if content_items:
        content_box = slide.shapes.add_textbox(Inches(0.75), Inches(y_offset), Inches(11.833), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(24)
            p.font.color.rgb = COLORS['text']
            p.space_after = Pt(12)

    return slide


def add_quote_slide(prs, quote):
    """Add a quote slide."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['dark']
    background.line.fill.background()

    # Quote
    quote_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10.333), Inches(3.5))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f'"{quote}"'
    p.font.size = Pt(32)
    p.font.italic = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_discussion_slide(prs, title, question, additional=None):
    """Add a discussion slide."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Background with gradient effect (simulated with shape)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0xFA, 0xF9, 0xF6)
    background.line.fill.background()

    # Accent bar at top
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLORS['primary']
    accent_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1), Inches(11.833), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['text']
    p.alignment = PP_ALIGN.CENTER

    # Question
    q_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.333), Inches(1.5))
    tf = q_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = question
    p.font.size = Pt(36)
    p.font.color.rgb = COLORS['primary']
    p.alignment = PP_ALIGN.CENTER

    if additional:
        add_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(1))
        tf = add_box.text_frame
        p = tf.paragraphs[0]
        p.text = additional
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['text']
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_two_column_slide(prs, title, left_items, right_items, left_header=None, right_header=None):
    """Add a two-column comparison slide."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0xFA, 0xF9, 0xF6)
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['text']

    y_start = 1.8

    # Left column header
    if left_header:
        lh_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.5), Inches(0.5))
        tf = lh_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_header
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        y_start = 2.1

    # Right column header
    if right_header:
        rh_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.5), Inches(0.5))
        tf = rh_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_header
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLORS['coral']

    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.75), Inches(y_start), Inches(5.5), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['text']
        p.space_after = Pt(10)

    # Right column content
    right_box = slide.shapes.add_textbox(Inches(7), Inches(y_start), Inches(5.5), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['text']
        p.space_after = Pt(10)

    return slide


def add_stats_slide(prs, title, stats):
    """Add a statistics slide with large numbers."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0xFA, 0xF9, 0xF6)
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['text']

    # Stats cards
    card_width = 3.5
    total_width = card_width * len(stats) + 0.5 * (len(stats) - 1)
    start_x = (13.333 - total_width) / 2

    colors = [COLORS['primary'], COLORS['coral'], COLORS['blue']]

    for i, (value, desc, source) in enumerate(stats):
        x = start_x + i * (card_width + 0.5)

        # Stat value
        val_box = slide.shapes.add_textbox(Inches(x), Inches(2.5), Inches(card_width), Inches(1.2))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = colors[i % len(colors)]
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(Inches(x), Inches(3.7), Inches(card_width), Inches(1.2))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['text']
        p.alignment = PP_ALIGN.CENTER

        # Source
        src_box = slide.shapes.add_textbox(Inches(x), Inches(5), Inches(card_width), Inches(0.5))
        tf = src_box.text_frame
        p = tf.paragraphs[0]
        p.text = source
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_capability_slide(prs, capability_name, definition, color=None):
    """Add a capability section title slide."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    bg_color = color if color else COLORS['primary']

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = bg_color
    background.line.fill.background()

    # Capability name
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = capability_name
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Definition
    def_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.333), Inches(1))
    tf = def_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = definition
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_story_slide(prs, title, story_text, impact=None):
    """Add a real story slide."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0xFA, 0xF9, 0xF6)
    background.line.fill.background()

    # Accent shape
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.8), Inches(0.1), Inches(2)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS['primary']
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['text']

    # Story text
    story_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(2))
    tf = story_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = story_text
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['text']

    # Impact
    if impact:
        impact_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(1))
        tf = impact_box.text_frame
        p = tf.paragraphs[0]
        p.text = impact
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_takeaways_slide(prs, title, takeaways):
    """Add a takeaways slide with numbered items."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0xFA, 0xF9, 0xF6)
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['text']

    # Takeaways
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.333), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(takeaways):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{i + 1}. {item}"
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['text']
        p.space_after = Pt(16)

    return slide


def add_transition_slide(prs, text, next_text=None):
    """Add a transition slide."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['dark']
    background.line.fill.background()

    # Main text
    text_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.333), Inches(1.5))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if next_text:
        next_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.333), Inches(1))
        tf = next_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = next_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.alignment = PP_ALIGN.CENTER

    return slide


def main():
    """Generate the PowerPoint presentation."""
    prs = create_presentation()

    # ===== PART I: THE LANDSCAPE =====

    # Slide 1: Title
    add_title_slide(prs, "Executive AI Briefing", "Understanding AI for Business Leaders", "FENLO AI")

    # Slide 2: What to Expect
    add_content_slide(prs, "What to Expect", [
        "Part I: The Landscape - Understanding AI fundamentals",
        "Part II: Industry Scope - How AI drives business value",
        "Part III: Opportunities - Exploring possibilities for you",
        "Part IV: Open Forum - Your questions answered",
    ], '"No question is too basic" - this is your space to learn')

    # Slide 3: Quick Pulse Check
    add_discussion_slide(prs, "Quick Pulse Check", 'What comes to mind when you hear "AI"?')

    # Slide 4: Section Title - AI Timeline
    add_section_title_slide(prs, "The AI Timeline", "A 70-Year Journey to This Moment")

    # Slide 5: 1950s-1960s
    add_content_slide(prs, "1950s-1960s: The Dream Begins", [
        '1950: Alan Turing asks "Can machines think?"',
        '1956: "Artificial Intelligence" coined at Dartmouth',
        'Early optimism: "Human-level AI in 20 years!"',
        "Reality: Computers too slow, memory too limited",
        "Fun fact: The 1956 Dartmouth workshop had just 10 attendees"
    ])

    # Slide 6: 1970s-1980s
    add_content_slide(prs, "1970s-1980s: The First AI Winter", [
        "Grand promises didn't materialize",
        "Funding dried up dramatically",
        'AI became a "dirty word" in academia',
        "But quiet progress continued in the shadows..."
    ])

    # Slide 7: 1990s-2000s
    add_content_slide(prs, "1990s-2000s: AI Goes Undercover", [
        "AI stopped trying to replicate human intelligence",
        "Instead, it solved specific problems really well",
        "And it quietly entered your life..."
    ])

    # Slide 8: AI You Didn't Know Was AI
    add_content_slide(prs, "AI You Didn't Know Was AI", [
        "Spam Filters - Keeping your inbox clean",
        'Netflix "You Might Like" - Recommendation engines',
        "Google Search Ranking - Finding what you need",
        "Fraud Detection - Protecting your bank account",
    ], 'AI was everywhere - we just didn\'t call it AI')

    # Slide 9: 2010s
    add_content_slide(prs, "2010s: The Deep Learning Revolution", [
        "Massive computing power (GPUs!) + huge data",
        "2012: AI beats humans at image recognition",
        "Tech giants invest billions",
        'Fun fact: The 2012 "AlexNet" moment used GPUs designed for video games'
    ])

    # Slide 10: Milestone Moments
    add_content_slide(prs, "Milestone Moments", [
        "2011: IBM Watson wins Jeopardy!",
        "2014: Alexa enters homes",
        "2016: AlphaGo defeats world champion"
    ])

    # Slide 11: 2020s - The AI Explosion
    add_stats_slide(prs, "2020s - The AI Explosion", [
        ("100M", "users in just 2 months", "ChatGPT Launch - Nov 2022"),
        ("2 months", "Fastest tech adoption in history", "vs 2.5 years for Instagram"),
        ("70 years", "of AI research coming together", "The culmination of decades")
    ])

    # Slide 12: Quote
    add_quote_slide(prs, "AI didn't appear suddenly - it's been building for 70 years. What's new is that it finally became good enough and accessible enough for everyone to use.")

    # Slide 13: Discussion
    add_discussion_slide(prs, "Discussion", "Does this timeline change how you think about AI?", 'Any "aha" moments?')

    # Slide 14: Section Title - AI You Already Use
    add_section_title_slide(prs, "AI You Already Use", "It's more familiar than you think")

    # Slide 15: On Your Phone
    add_content_slide(prs, "On Your Phone", [
        "Face Unlock (Vision AI)",
        "Autocorrect & Predictive Text",
        'Photo Organization ("Show me beach photos")',
        "Siri, Google Assistant, Alexa"
    ])

    # Slide 16: At Work
    add_content_slide(prs, "At Work", [
        "Spam Filtering (you never see 90% of spam)",
        "Meeting Transcription (Zoom, Teams)",
        "Smart Document Search",
        "Calendar Scheduling Suggestions"
    ])

    # Slide 17: In Daily Life
    add_content_slide(prs, "In Daily Life", [
        "Navigation & Traffic Predictions",
        "Netflix, Spotify Recommendations",
        'Amazon "Customers also bought..."',
        "Bank Fraud Alerts"
    ])

    # Slide 18: Quote
    add_quote_slide(prs, "You've been using AI for years. What's new is that now you can direct it yourself, instead of it working invisibly in the background.")

    # Slide 19: Section Title - Breakthroughs
    add_section_title_slide(prs, "The Breakthroughs", "Understanding Different Types of AI")

    # Slide 20: Deep Learning - The Foundation
    add_two_column_slide(prs, "Deep Learning - The Foundation",
        ["Humans write rules", "Computers follow blindly", "Limited to what programmers anticipate"],
        ["Show thousands of examples", "Computer learns patterns", "Discovers insights humans miss"],
        "Traditional Programming", "Deep Learning")

    # Slide 21: Quote
    add_quote_slide(prs, "Instead of teaching 'a cat has pointy ears, whiskers, and fur,' you show 10,000 pictures until it figures out what makes a cat a cat.")

    # Slide 22: Vision Models
    add_content_slide(prs, "Vision Models: Machines That See", [
        "Recognize Faces",
        "Read Documents & Handwriting",
        "Detect Manufacturing Defects",
        "Analyze Medical Scans",
        "Identify Objects in Video",
        "Read License Plates"
    ])

    # Slide 23: Vision Models - Business Applications
    add_content_slide(prs, "Vision Models - Business Applications", [
        "Quality control - catches defects humans miss",
        "Document processing - extracts data from PDFs",
        "Security - face recognition access",
        "Insurance - assess damage from photos"
    ])

    # Slide 24: Language Models
    add_content_slide(prs, "Language Models: Machines That Read & Write", [
        '2017: "Transformers" breakthrough at Google',
        "AI finally understands context and meaning",
        "Not just matching keywords anymore - actually gets it"
    ])

    # Slide 25: What Language Models Do
    add_content_slide(prs, "What Language Models Do", [
        "Answer Complex Questions",
        "Write Any Type of Content",
        "Summarize Long Documents",
        "Translate 100+ Languages",
        "Extract Key Information",
        "Analyze Sentiment & Tone"
    ])

    # Slide 26: Large Language Models
    add_content_slide(prs, "Large Language Models (LLMs)", [
        "Trained on most of the internet's text",
        "Handle almost any text-based task",
        "No coding or technical skills required",
        "Understand context, nuance, even humor"
    ], "The Current Wave")

    # Slide 27: The Key Players
    add_content_slide(prs, "The Key Players", [
        "OpenAI - ChatGPT / GPT-4",
        "Anthropic - Claude",
        "Google - Gemini",
        "Microsoft - Copilot",
        "Meta - Llama (Open Source)"
    ])

    # Slide 28: Which One Should You Use
    add_two_column_slide(prs, "Which One Should You Use?",
        ["ChatGPT: General-purpose, creative writing, large plugin ecosystem",
         "Claude: Long documents (200k+ tokens), nuanced responses, enterprise security"],
        ["Gemini: Google Workspace integration, multimodal, free tier",
         "Copilot: Microsoft integration, great for Office users"],
        "Options", "More Options")

    # Slide 29: LLM Limitations
    add_content_slide(prs, "LLM Limitations - Important to Know", [
        "Hallucinations: Can be confidently wrong - always verify facts",
        "Knowledge cutoff: Doesn't know recent events",
        "Garbage in, garbage out: Quality depends on your instructions"
    ])

    # Slide 30: AI Agents
    add_two_column_slide(prs, "AI Agents: Machines That Take Action",
        ["You ask a question", "It gives an answer", "You do the work"],
        ["You give a goal", "It figures out steps", "It takes action for you"],
        "Regular AI", "AI Agent")

    # Slide 31: Agent Examples
    add_content_slide(prs, "Agent Examples", [
        '"Research our top 5 competitors" → Browses websites, reads reviews, compiles report',
        '"Schedule a team meeting next week" → Checks calendars, finds time, sends invites',
        '"Book me a flight to NYC under $500" → Searches options, compares prices, recommends best'
    ])

    # Slide 32: Quote
    add_quote_slide(prs, "Regular AI is like a calculator - you push buttons and get answers. An agent is like a personal assistant who can run errands for you.")

    # Slide 33: Generative AI
    add_content_slide(prs, "Generative AI: Machines That Create", [
        "Not just analyzing existing content...",
        "Creating brand new text, images, code, audio, video",
        'This is the "magic" people are excited about'
    ])

    # Slide 34: What GenAI Can Create
    add_content_slide(prs, "What Generative AI Can Create", [
        "Text: Emails, reports, articles, scripts",
        "Images: Marketing visuals, mockups, art",
        "Code: Programs, websites, scripts",
        "Audio: Voiceovers, music, podcasts",
        "Video: Clips, animations, ads"
    ])

    # Slide 35: Quote
    add_quote_slide(prs, "Generative AI amplifies human creativity - it doesn't replace the need for judgment, taste, and oversight.")

    # Slide 36: Key Terms
    add_content_slide(prs, "Key Terms - Plain English", [
        "Prompt: The instruction or question you give to AI",
        "Hallucination: When AI confidently makes up facts",
        "Training: How AI learns from millions of examples",
        "Model: The AI 'brain' (GPT-4, Claude, Gemini)",
        "Context Window: How much AI can 'remember' at once",
        "Fine-tuning: Customizing AI for specific tasks"
    ])

    # Slide 37: Discussion
    add_discussion_slide(prs, "Discussion", "Which AI type seems most relevant to your business?", "What surprised you most? What concerns do you have?")

    # Slide 38: Part I Takeaways
    add_takeaways_slide(prs, "Part I Takeaways", [
        "AI has been building for 70 years - this isn't hype",
        "You already use AI daily - it's familiar",
        "Different types serve different purposes",
        "We're in a pivotal moment - accessible to all",
        "Human judgment remains essential"
    ])

    # Slide 39: Transition
    add_transition_slide(prs, "Now that we understand what AI is and how we got here...", "Let's see how businesses are actually using it to create value.")

    # ===== PART II: INDUSTRY SCOPE =====

    # Slide 40: Section Title
    add_section_title_slide(prs, "Industry Scope", "How AI Drives Real Business Value")

    # Slide 41: Reality Check
    add_content_slide(prs, "The Reality Check", [
        "AI amplifies people, doesn't replace them",
        "Winners start with problems, not technology",
        "Fortune 500s AND 10-person teams benefit"
    ])

    # Slide 42: The Numbers Are Real
    add_stats_slide(prs, "The Numbers Are Real", [
        ("40%", "productivity gain in AI-augmented workflows", "McKinsey 2024"),
        ("72%", "of executives say AI is a top priority", "Deloitte 2024"),
        ("$15.7T", "potential AI contribution to global economy by 2030", "PwC Global AI Study")
    ])

    # Slide 43: 6 AI Capabilities
    add_content_slide(prs, "6 AI Capabilities", [
        "AUTOMATE - Handle repetitive tasks (Start here!)",
        "ANALYZE - Find patterns and insights",
        "CREATE - Generate content",
        "ASSIST - Help humans work",
        "PREDICT - Forecast the future",
        "PERSONALIZE - Tailor to individuals"
    ])

    # Slide 44: AUTOMATE Title
    add_capability_slide(prs, "AUTOMATE", "AI handles repetitive tasks so your team can focus on what matters", COLORS['primary'])

    # Slide 45: AUTOMATE Before & After
    add_two_column_slide(prs, "AUTOMATE - Before & After",
        ["Manual invoice entry - hours per week",
         "Owner reads every email",
         "Phone tag for scheduling",
         "Data entry from paper forms"],
        ["Auto-extracted and filed in seconds",
         "Auto-categorized, urgent items flagged",
         "Self-service AI booking assistant",
         "Scan → auto-captured → database"],
        "Before AI", "After AI")

    # Slide 46: AUTOMATE Enterprise
    add_content_slide(prs, "AUTOMATE - Enterprise Scale", [
        "Claims processing: Weeks of manual review → Same-day decisions",
        "Employee onboarding: HR manually provisions → Automated workflows",
        "Compliance monitoring: Quarterly audits → Continuous automated checks"
    ])

    # Slide 47: AUTOMATE Story
    add_story_slide(prs, "Real Story", '"A 15-person accounting firm automated invoice processing. What took 2 hours daily now takes 10 minutes."', "40+ hours/month back for client work")

    # Slide 48: ANALYZE Title
    add_capability_slide(prs, "ANALYZE", "AI finds patterns and insights hidden in your data", COLORS['blue'])

    # Slide 49: ANALYZE Examples
    add_content_slide(prs, "ANALYZE - Ask AI Questions", [
        '"Which customers are likely to buy again this quarter?"',
        '"What themes are appearing in our customer reviews?"',
        '"Flag any unusual expenses in this report"',
        '"Summarize key trends in plain English"'
    ])

    # Slide 50: ANALYZE Enterprise
    add_content_slide(prs, "ANALYZE - Enterprise Power", [
        "Supply chain: Predict disruptions before they happen",
        "Churn prediction: Identify at-risk customers proactively",
        "Forecasting: Multiple scenarios in minutes, not days"
    ])

    # Slide 51: ANALYZE Story
    add_story_slide(prs, "Real Story", '"A retail chain found that weather 5 days out was the best predictor of inventory needs."', "Something no human analyst had ever noticed")

    # Slide 52: CREATE Title
    add_capability_slide(prs, "CREATE", "AI generates content from your direction - first drafts in seconds", COLORS['coral'])

    # Slide 53: CREATE Examples
    add_two_column_slide(prs, "CREATE - Content at Scale",
        ["Marketing emails & newsletters",
         "Social media posts (week in minutes!)",
         "Product descriptions",
         "Proposals & quotes"],
        ["Training materials & assessments",
         "Company-wide communications",
         "Technical documentation",
         "Localization (50+ languages)"],
        "SMB Use Cases", "Enterprise Use Cases")

    # Slide 54: CREATE Quote
    add_quote_slide(prs, "AI creates first drafts, not final products. Human review ensures quality, accuracy, and your unique brand voice.")

    # Slide 55: ASSIST Title
    add_capability_slide(prs, "ASSIST", "AI as a knowledgeable helper, available 24/7", COLORS['purple'])

    # Slide 56: ASSIST Examples
    add_two_column_slide(prs, "ASSIST - Your AI Teammate",
        ["Customer support chatbot",
         "Internal FAQ assistant",
         "Sales call prep",
         "Meeting notes & action items"],
        ["IT help desk (auto-resolve tickets)",
         "HR policy questions 24/7",
         "Legal document research",
         "Engineering code copilot"],
        "SMB Applications", "Enterprise Applications")

    # Slide 57: ASSIST Story
    add_story_slide(prs, "Real Story", '"A law firm trained AI on their document library. Junior associates now get answers in seconds that used to require bothering senior partners."', "Senior time freed up. Juniors empowered. Everyone wins.")

    # Slide 58: PREDICT Title
    add_capability_slide(prs, "PREDICT", "AI forecasts future outcomes so you can prepare, not react", COLORS['teal'])

    # Slide 59: PREDICT Examples
    add_two_column_slide(prs, "PREDICT - See the Future",
        ["Cash flow forecasting",
         "Demand planning",
         "Lead scoring",
         "Equipment maintenance"],
        ["Know when money will be tight",
         "Stock the right inventory",
         "Focus on best prospects",
         "Fix it before it breaks"],
        "What AI Predicts", "Business Impact")

    # Slide 60: PREDICT Story
    add_story_slide(prs, "Real Story", '"A manufacturing company predicts equipment failures 2 weeks in advance."', "Unplanned downtime dropped 70%")

    # Slide 61: PERSONALIZE Title
    add_capability_slide(prs, "PERSONALIZE", "AI tailors every experience to each individual", COLORS['gold'])

    # Slide 62: PERSONALIZE Examples
    add_two_column_slide(prs, "PERSONALIZE - One Size Fits One",
        ["Segmented email campaigns",
         '"Customers like you bought..."',
         "Dynamic pricing",
         "Optimal send timing"],
        ["Website personalization by visitor",
         "Adaptive employee training paths",
         "Next-best-action in sales",
         "Multi-channel journey orchestration"],
        "SMB Applications", "Enterprise Applications")

    # Slide 63: PERSONALIZE Story
    add_story_slide(prs, "Real Story", '"An e-commerce SMB personalized email subject lines with AI."', "35% higher open rates")

    # Slide 64: Value Stack
    add_content_slide(prs, "The AI Value Stack", [
        "PERSONALIZE - Tailor to individuals",
        "PREDICT - Forecast the future",
        "ASSIST - Help humans work",
        "CREATE - Generate content",
        "ANALYZE - Find insights",
        "AUTOMATE - Handle repetitive tasks ← Start here!"
    ])

    # Slide 65: Quote
    add_quote_slide(prs, "You don't need all six. Start with one capability that solves a real problem. Expand from there.")

    # Slide 66: Discussion
    add_discussion_slide(prs, "Discussion", "Which capability would make the biggest difference for you?", "Did any examples remind you of current challenges?")

    # Slide 67: Part II Takeaways
    add_takeaways_slide(prs, "Part II Takeaways", [
        "6 core ways AI creates business value",
        "Size doesn't matter - SMBs and enterprises both win",
        "Start with the problem, not the technology",
        "Humans remain central - AI is a tool",
        "Results are real and measurable"
    ])

    # Slide 68: Transition
    add_transition_slide(prs, "You've seen what's possible across industries.", "Now let's focus on your organization.\n\nPart III: Opportunity Exploration →")

    # Save
    output_path = "/Users/mohammadshoaib/Codes/Executive Brief/slides/Executive_AI_Briefing.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
